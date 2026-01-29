import os
from dotenv import load_dotenv
load_dotenv()

from flask import Flask, render_template, request, redirect, url_for, send_from_directory, flash, after_this_request, session
import pythoncom
from flask_compress import Compress
from werkzeug.utils import secure_filename
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from pdf2docx import Converter
from docx2pdf import convert
from pdf2image import convert_from_path
from PIL import Image
import img2pdf
from openpyxl import load_workbook
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
import io
import tempfile
import pytesseract
import openai
from transformers import pipeline
import torch
import time

# Configure the path to the Tesseract executable
pytesseract.pytesseract.tesseract_cmd = r'C:/Program Files/Tesseract-OCR/tesseract.exe'  # Update this path if your installation is different

import time
from flask import g

app = Flask(__name__, static_url_path='/static', static_folder='static')
Compress(app)  # Enable gzip compression

app.config['SECRET_KEY'] = 'your-secret-key-here'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['PROCESSED_FOLDER'] = 'processed'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'doc', 'jpg', 'jpeg', 'png', 'bmp', 'gif', 'tiff', 'webp', 'xlsx', 'pptx'}

# Limit max upload size to 16MB (adjust as needed)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16 MB

# Max upload size in MB for display
app.config['MAX_UPLOAD_SIZE_MB'] = 16

@app.errorhandler(413)
def request_entity_too_large(error):
    return "File is too large. Maximum allowed size is 16 MB.", 413

# Add Ghostscript path config here - update this path if your Ghostscript is installed elsewhere
app.config['GHOSTSCRIPT_PATH'] = r"C:\\Program Files\\gs\\gs10.05.1\\bin\\gswin64c.exe"

# Ensure directories exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)

@app.before_request
def start_timer():
    g.start_time = time.time()

@app.after_request
def log_request_performance(response):
    if hasattr(g, 'start_time'):
        duration = time.time() - g.start_time
        # Only log requests that take longer than 0.1 seconds to reduce log noise and overhead
        if duration > 0.1:
            endpoint = request.endpoint
            method = request.method
            status_code = response.status_code
            app.logger.info(f"{method} {endpoint} {status_code} completed in {duration:.4f} seconds")
    # Add cache control headers for static files
    if request.path.startswith('/static/'):
        response.headers['Cache-Control'] = 'public, max-age=31536000'
    return response

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def allowed_pdf_file(file):
    # Check extension and MIME type for PDF files
    if not file:
        return False
    filename = file.filename
    if '.' not in filename or filename.rsplit('.', 1)[1].lower() != 'pdf':
        return False
    if file.mimetype != 'application/pdf':
        return False
    return True

from dotenv import load_dotenv
load_dotenv(override=True)  # Ensure .env variables override existing environment variables

import os
import logging

openai_api_key = os.getenv("OPENAI_API_KEY")
if not openai_api_key:
    logging.error("OPENAI_API_KEY environment variable is not set or empty.")
else:
    logging.info(f"OPENAI_API_KEY loaded successfully: {openai_api_key[:4]}{'*'*(len(openai_api_key)-8)}{openai_api_key[-4:]}")

openai.api_key = openai_api_key

# Initialize Hugging Face summarization pipeline
summarizer = None

def get_summarizer():
    global summarizer
    if summarizer is None:
        try:
            summarizer = pipeline("summarization", model="facebook/bart-large-cnn", device=0 if torch.cuda.is_available() else -1)
            app.logger.info("Hugging Face summarization pipeline loaded successfully.")
        except Exception as e:
            summarizer = None
            app.logger.error(f"Failed to load Hugging Face summarization pipeline: {e}")
    return summarizer

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/debug-api-key')
def debug_api_key():
    key = os.getenv("OPENAI_API_KEY", "")
    masked_key = key[:4] + "*" * (len(key) - 8) + key[-4:] if len(key) > 8 else "Not Set"
    return f"Current OPENAI_API_KEY: {masked_key}"

@app.route('/summarize', methods=['GET', 'POST'])
def summarize():
    import logging
    logging.info(f"OpenAI API Key: {openai.api_key}")
    summary = None
    if request.method == 'POST':
        text = ''
        # Check if a file was uploaded
        if 'file' in request.files:
            file = request.files['file']
            if file and file.filename != '' and allowed_pdf_file(file):
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                # Extract text from file based on extension
                ext = filename.rsplit('.', 1)[1].lower()
                if ext == 'pdf':
                    try:
                        reader = PdfReader(filepath)
                        for page in reader.pages:
                            text += page.extract_text() or ''
                    except Exception as e:
                        flash('Failed to extract text from PDF.')
                        return redirect(request.url)
            elif file and file.filename != '' and allowed_file(file.filename):
                # For other allowed file types (docx, doc)
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                ext = filename.rsplit('.', 1)[1].lower()
                if ext in ['docx', 'doc']:
                    try:
                        from docx import Document
                        doc = Document(filepath)
                        for para in doc.paragraphs:
                            text += para.text + '\n'
                    except Exception as e:
                        flash('Failed to extract text from Word document.')
                        return redirect(request.url)
                else:
                    flash('Unsupported file type for summarization.')
                    return redirect(request.url)
            else:
                # No valid file uploaded, but check if text input is provided
                text = request.form.get('text', '')
                if not text.strip():
                    flash('Please provide text or upload a document for summarization.')
                    return redirect(request.url)
        else:
            # If no file, check for text input
            text = request.form.get('text', '')
        if not text.strip():
            flash('Please provide text or upload a document for summarization.')
            return redirect(request.url)
        # Call OpenAI API to generate summary
        try:
            response = openai.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that summarizes academic documents concisely."},
                    {"role": "user", "content": f"Please provide a concise summary of the following text:\n\n{text}"}
                ],
                max_tokens=300,
                temperature=0.5,
            )
            summary = response.choices[0].message.content
        except Exception as e:
            import traceback
            logging.error(f"OpenAI API call failed: {e}")
            logging.error(traceback.format_exc())
            error_str = str(e).lower()
            logging.error(f"Exception type: {type(e)}")
            if "insufficient_quota" in error_str or "429" in error_str or "quota" in error_str:
                flash('OpenAI API quota exceeded. Using local summarization fallback.')
                local_summarizer = get_summarizer()
                if local_summarizer:
                    try:
                        # Use Hugging Face summarizer as fallback
                        # Ensure text length is within model limits (e.g., truncate if too long)
                        max_input_length = 1024  # typical max tokens for bart-large-cnn
                        truncated_text = text[:max_input_length]
                        hf_summary = local_summarizer(truncated_text, max_length=130, min_length=30, do_sample=False)
                        summary = hf_summary[0]['summary_text']
                    except Exception as hf_e:
                        app.logger.error(f"Hugging Face summarization failed: {hf_e}")
                        flash('Local summarization also failed.')
                        summary = None
                else:
                    flash('Local summarization model not available.')
                    summary = None
            elif "invalid_api_key" in error_str or "unauthorized" in error_str:
                flash('Invalid OpenAI API key. Please check your API key configuration.')
                summary = None
            else:
                flash(f'Failed to generate summary using OpenAI API. Error: {e}')
                summary = None
        return render_template('summarize.html', summary=summary)
    else:
        # For GET request, just render the summarize page without summary
        return render_template('summarize.html', summary=None)

# Merge PDFs
@app.route('/merge', methods=['GET', 'POST'])
def merge_pdfs():
    if request.method == 'POST':
        files = request.files.getlist('files')
        if len(files) < 2:
            flash('Please select at least 2 PDF files to merge.')
            return redirect(request.url)
        
        merger = PdfMerger()
        filenames = []
        
        for file in files:
            if file and allowed_pdf_file(file):
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                filenames.append(filename)
                merger.append(filepath)
        
        if len(filenames) < 2:
            flash('Please select at least 2 valid PDF files.')
            return redirect(request.url)
        
        output_filename = 'merged.pdf'
        output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
        merger.write(output_path)
        merger.close()
        
        return redirect(url_for('download_file', filename=output_filename))
    
    return render_template('partials/merge.html')

# PDF to Word
@app.route('/pdf-to-word', methods=['GET', 'POST'])
def pdf_to_word():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_pdf_file(file):
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            output_filename = filename.rsplit('.', 1)[0] + '.docx'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            cv = Converter(pdf_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/convert.html', action='pdf_to_word', title='PDF to Word')

# Word to PDF
@app.route('/word-to-pdf', methods=['GET', 'POST'])
def word_to_pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() in ['docx', 'doc']:
            filename = secure_filename(file.filename)
            word_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(word_path)
            
            output_filename = filename.rsplit('.', 1)[0] + '.pdf'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            try:
                pythoncom.CoInitialize()
                convert(word_path, output_path)
            finally:
                pythoncom.CoUninitialize()
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid Word document (DOCX or DOC).')
            return redirect(request.url)
    
    return render_template('partials/convert.html', action='word_to_pdf', title='Word to PDF')

# Split PDF
@app.route('/split-pdf', methods=['GET', 'POST'])
def split_pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_pdf_file(file):
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            try:
                start_page = int(request.form.get('start_page', 1))
                end_page = int(request.form.get('end_page', 1))
            except ValueError:
                flash('Please enter valid page numbers.')
                return redirect(request.url)
            
            reader = PdfReader(pdf_path)
            writer = PdfWriter()
            
            if start_page < 1 or end_page > len(reader.pages) or start_page > end_page:
                flash('Invalid page range.')
                return redirect(request.url)
            
            for page_num in range(start_page - 1, end_page):
                writer.add_page(reader.pages[page_num])
            
            output_filename = f'split_{filename}'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            with open(output_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/split.html')

import subprocess

import platform
import shutil

import os

def compress_pdf_alternative(input_path, output_path):
    """
    Alternative PDF compression without Ghostscript.
    This method re-writes the PDF using PyPDF2 which may reduce size slightly.
    """
    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        with open(output_path, 'wb') as f_out:
            writer.write(f_out)
        return True
    except Exception as e:
        app.logger.error(f"Alternative PDF compression failed: {e}")
        return False

def compress_pdf_ghostscript(input_path, output_path):
    """
    Compress PDF using Ghostscript command line.
    """
    import sys
    # Print PATH environment variable for debugging
    app.logger.info(f"Python executable: {sys.executable}")
    app.logger.info(f"PATH environment variable: {os.environ.get('PATH')}")

    # Use configured Ghostscript path if set
    gs_executable = app.config.get("GHOSTSCRIPT_PATH", None)
    if gs_executable and not os.path.isfile(gs_executable):
        app.logger.error(f"Configured Ghostscript path does not exist: {gs_executable}")
        return False

    if not gs_executable:
        # Determine Ghostscript executable based on OS
        gs_executable = "gs"
        system = platform.system()
        if system == "Windows":
            # Try common Ghostscript executable names on Windows
            possible_names = ["gswin64c.exe", "gswin32c.exe"]
            for name in possible_names:
                if shutil.which(name):
                    gs_executable = name
                    break
        else:
            # For Linux/Mac, assume 'gs' is in PATH
            if shutil.which("gs"):
                gs_executable = "gs"
            else:
                gs_executable = None

    if not gs_executable:
        app.logger.error("Ghostscript executable not found in system PATH or configured path.")
        return False

    gs_command = [
        gs_executable,
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        "-dPDFSETTINGS=/ebook",
        "-dNOPAUSE",
        "-dQUIET",
        "-dBATCH",
        f"-sOutputFile={output_path}",
        input_path
    ]
    try:
        subprocess.run(gs_command, check=True)
        return True
    except Exception as e:
        app.logger.error(f"Ghostscript compression failed: {e}")
        return False

# Compress PDF
@app.route('/compress-pdf', methods=['GET', 'POST'])
def compress_pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_pdf_file(file):
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            output_filename = f'compressed_{filename}'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            # Use alternative compression method without Ghostscript
            success = compress_pdf_alternative(pdf_path, output_path)
            if not success:
                flash('PDF compression failed.')
                return redirect(request.url)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/compress.html')

# PDF to PowerPoint
@app.route('/pdf-to-ppt', methods=['GET', 'POST'])
def pdf_to_ppt():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() == 'pdf':
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            try:
                # Removed hardcoded poppler_path to rely on system PATH environment variable
                images = convert_from_path(pdf_path)
            except Exception as e:
                flash('Failed to convert PDF to PowerPoint. Please ensure Poppler is installed and added to your system PATH.')
                return redirect(request.url)
            
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            
            for img in images:
                temp_img_path = os.path.join(tempfile.gettempdir(), f'temp_{len(prs.slides)}.jpg')
                img.save(temp_img_path, 'JPEG')
                
                slide = prs.slides.add_slide(blank_slide_layout)
                left = top = 0
                slide.shapes.add_picture(temp_img_path, left, top, prs.slide_width, prs.slide_height)
                
                os.remove(temp_img_path)
            
            output_filename = filename.rsplit('.', 1)[0] + '.pptx'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            prs.save(output_path)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/convert.html', action='pdf_to_ppt', title='PDF to PowerPoint')

# Excel to PDF
@app.route('/excel-to-pdf', methods=['GET', 'POST'])
def excel_to_pdf():
    import logging
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() in ['xlsx', 'xls']:
            filename = secure_filename(file.filename)
            excel_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            try:
                file.save(excel_path)
            except Exception as e:
                logging.error(f"Failed to save uploaded Excel file: {e}")
                flash('Failed to save the uploaded file. Please try again.')
                return redirect(request.url)
            
            try:
                wb = load_workbook(excel_path)
                ws = wb.active
            except Exception as e:
                logging.error(f"Failed to load Excel workbook: {e}")
                flash('Failed to read the Excel file. Please ensure it is a valid Excel document.')
                return redirect(request.url)
            
            output_filename = filename.rsplit('.', 1)[0] + '.pdf'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            try:
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                
                y = height - 40
                for row in ws.iter_rows(values_only=True):
                    x = 40
                    for cell in row:
                        c.drawString(x, y, str(cell) if cell is not None else '')
                        x += 100
                    y -= 20
                    if y < 40:
                        c.showPage()
                        y = height - 40
                
                c.save()
            except Exception as e:
                logging.error(f"Failed to generate PDF from Excel: {e}")
                flash('Failed to generate PDF from the Excel file.')
                return redirect(request.url)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid Excel file (XLSX or XLS).')
            return redirect(request.url)
    
    return render_template('partials/convert.html', action='excel_to_pdf', title='Excel to PDF')


# PDF to JPG
@app.route('/pdf-to-jpg', methods=['GET', 'POST'])
def pdf_to_jpg():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() == 'pdf':
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            poppler_path = r"C:\\poppler\\bin"  # Update this path to your actual poppler bin path
            try:
                images = convert_from_path(pdf_path)
            except Exception as e:
                flash('Failed to convert PDF to JPG. Please ensure Poppler is installed and added to your system PATH.')
                return redirect(request.url)
            
            if len(images) == 1:
                output_filename = filename.rsplit('.', 1)[0] + '.jpg'
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                images[0].save(output_path, 'JPEG')
            else:
                import zipfile
                output_filename = filename.rsplit('.', 1)[0] + '.zip'
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                
                with zipfile.ZipFile(output_path, 'w') as zipf:
                    for i, image in enumerate(images):
                        img_path = os.path.join(tempfile.gettempdir(), f'page_{i+1}.jpg')
                        image.save(img_path, 'JPEG')
                        zipf.write(img_path, f'page_{i+1}.jpg')
                        os.remove(img_path)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/convert.html', action='pdf_to_jpg', title='PDF to JPG')

# JPG to PDF
@app.route('/jpg-to-pdf', methods=['GET', 'POST'])
def jpg_to_pdf():
    if request.method == 'POST':
        files = request.files.getlist('files')
        if not files or len(files) == 0:
            flash('Please select at least one image file.')
            return redirect(request.url)
        
        images = []
        for file in files:
            if file and file.filename != '' and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() in ['jpg', 'jpeg', 'png']:
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                images.append(filepath)
        
        if not images:
            flash('Please select valid image files (JPG, JPEG, PNG).')
            return redirect(request.url)
        
        output_filename = 'converted.pdf'
        output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
        
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(images))
        
        for img_path in images:
            os.remove(img_path)
        
        return redirect(url_for('download_file', filename=output_filename))
    
    return render_template('partials/jpg_to_pdf.html')

# Add Watermark
from reportlab.pdfgen import canvas as pdf_canvas
from reportlab.lib.pagesizes import letter
from PyPDF2 import PdfReader, PdfWriter

@app.route('/add-watermark', methods=['GET', 'POST'])
def add_watermark():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('Please select a PDF file.')
            return redirect(request.url)
        
        pdf_file = request.files['file']
        watermark_text = request.form.get('watermark_text', '').strip()
        
        if pdf_file.filename == '':
            flash('Please select a PDF file.')
            return redirect(request.url)
        
        if not watermark_text:
            flash('Please enter watermark text.')
            return redirect(request.url)
        
        if pdf_file and allowed_file(pdf_file.filename) and pdf_file.filename.rsplit('.', 1)[1].lower() == 'pdf':
            pdf_filename = secure_filename(pdf_file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)
            pdf_file.save(pdf_path)
            
            reader = PdfReader(pdf_path)
            writer = PdfWriter()
            
            # Create watermark PDF with text matching each page size
            watermark_pages = []
            for page in reader.pages:
                width = float(page.mediabox.width)
                height = float(page.mediabox.height)
                packet = io.BytesIO()
                c = pdf_canvas.Canvas(packet, pagesize=(width, height))
                c.setFont("Helvetica", 40)
                c.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.3)
                c.saveState()
                c.translate(width/2, height/2)
                c.rotate(45)
                c.drawCentredString(0, 0, watermark_text)
                c.restoreState()
                c.save()
                packet.seek(0)
                watermark_pdf = PdfReader(packet)
                watermark_pages.append(watermark_pdf.pages[0])
            
            for i, page in enumerate(reader.pages):
                page.merge_page(watermark_pages[i])
                writer.add_page(page)
            
            output_filename = f'watermarked_{pdf_filename}'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            with open(output_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            os.remove(pdf_path)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/watermark.html')

# Rotate PDF
@app.route('/rotate-pdf', methods=['GET', 'POST'])
def rotate_pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() == 'pdf':
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            try:
                angle = int(request.form.get('angle', 0))
                if angle not in [0, 90, 180, 270]:
                    flash('Rotation must be 0, 90, 180, or 270 degrees.')
                    return redirect(request.url)
            except ValueError:
                flash('Invalid rotation value.')
                return redirect(request.url)
            
            reader = PdfReader(pdf_path)
            writer = PdfWriter()
            
            for page in reader.pages:
                if angle == 90:
                    page.rotate(90)
                elif angle == 180:
                    page.rotate(180)
                elif angle == 270:
                    page.rotate(-90)
                # angle == 0 means no rotation
                
                writer.add_page(page)
            
            output_filename = f'rotated_{filename}'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            with open(output_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/rotate.html')

# Protect PDF
@app.route('/protect-pdf', methods=['GET', 'POST'])
def protect_pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        password = request.form.get('password', '')
        if not password:
            flash('Please enter a password.')
            return redirect(request.url)
        
        if file and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() == 'pdf':
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            reader = PdfReader(pdf_path)
            writer = PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            writer.encrypt(password)
            
            output_filename = f'protected_{filename}'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            with open(output_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/protect.html')

# Unlock PDF
@app.route('/unlock-pdf', methods=['GET', 'POST'])
def unlock_pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        password = request.form.get('password', '')
        
        if file and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() == 'pdf':
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            try:
                reader = PdfReader(pdf_path)
                if reader.is_encrypted:
                    if not password:
                        flash('This PDF is encrypted. Please provide the password.')
                        return redirect(request.url)
                    reader.decrypt(password)
                
                writer = PdfWriter()
                
                for page in reader.pages:
                    writer.add_page(page)
                
                output_filename = f'unlocked_{filename}'
                output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
                
                with open(output_path, 'wb') as output_pdf:
                    writer.write(output_pdf)
                
                return redirect(url_for('download_file', filename=output_filename))
            except Exception as e:
                flash('Failed to unlock PDF. Incorrect password or corrupt file.')
                return redirect(request.url)
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/unlock.html')

# Organize PDF
@app.route('/organize-pdf', methods=['GET', 'POST'])
def organize_pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() == 'pdf':
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            page_order = request.form.get('page_order', '')
            if not page_order:
                flash('Please specify the page order.')
                return redirect(request.url)
            
            try:
                order = [int(x.strip()) for x in page_order.split(',')]
            except ValueError:
                flash('Invalid page order format. Use comma-separated numbers (e.g., "1,3,2").')
                return redirect(request.url)
            
            reader = PdfReader(pdf_path)
            if any(p < 1 or p > len(reader.pages) for p in order):
                flash(f'Page numbers must be between 1 and {len(reader.pages)}.')
                return redirect(request.url)
            
            writer = PdfWriter()
            
            for p in order:
                writer.add_page(reader.pages[p-1])
            
            output_filename = f'organized_{filename}'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            with open(output_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/organize.html')

# Delete Pages from PDF
@app.route('/delete-pages', methods=['GET', 'POST'])
def delete_pages():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected.')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected.')
            return redirect(request.url)
        
        if file and allowed_file(file.filename) and file.filename.rsplit('.', 1)[1].lower() == 'pdf':
            filename = secure_filename(file.filename)
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(pdf_path)
            
            pages_to_delete = request.form.get('pages_to_delete', '')
            if not pages_to_delete:
                flash('Please specify which pages to delete.')
                return redirect(request.url)
            
            try:
                delete_list = [int(x.strip()) for x in pages_to_delete.split(',')]
            except ValueError:
                flash('Invalid page numbers format. Use comma-separated numbers (e.g., "2,5").')
                return redirect(request.url)
            
            reader = PdfReader(pdf_path)
            if any(p < 1 or p > len(reader.pages) for p in delete_list):
                flash(f'Page numbers must be between 1 and {len(reader.pages)}.')
                return redirect(request.url)
            
            writer = PdfWriter()
            
            for i in range(len(reader.pages)):
                if (i+1) not in delete_list:
                    writer.add_page(reader.pages[i])
            
            output_filename = f'deleted_pages_{filename}'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            
            with open(output_path, 'wb') as output_pdf:
                writer.write(output_pdf)
            
            return redirect(url_for('download_file', filename=output_filename))
        else:
            flash('Please select a valid PDF file.')
            return redirect(request.url)
    
    return render_template('partials/delete_pages.html')

# Download file
import logging

import threading
import time

def delete_files_after_delay(processed_filename, delay=10):
    time.sleep(delay)
    try:
        # Delete processed file
        processed_path = os.path.join(app.config['PROCESSED_FOLDER'], processed_filename)
        if os.path.exists(processed_path):
            os.remove(processed_path)
            app.logger.info(f"Deleted processed file: {processed_path}")
        else:
            app.logger.warning(f"Processed file not found for deletion: {processed_path}")

        # Attempt to delete corresponding uploaded file
        base_name = os.path.splitext(processed_filename)[0]
        # Search for matching file in uploads folder with any allowed extension
        for ext in app.config['ALLOWED_EXTENSIONS']:
            upload_filename = base_name + '.' + ext
            upload_path = os.path.join(app.config['UPLOAD_FOLDER'], upload_filename)
            if os.path.exists(upload_path):
                os.remove(upload_path)
                app.logger.info(f"Deleted uploaded file: {upload_path}")
                break
    except Exception as e:
        app.logger.error(f"Error deleting files after delay: {e}", exc_info=True)

@app.route('/download/<filename>')
def download_file(filename):
    try:
        response = send_from_directory(app.config['PROCESSED_FOLDER'], filename, as_attachment=True)
        # Start background thread to delete files after 10 seconds
        threading.Thread(target=delete_files_after_delay, args=(filename, 10), daemon=True).start()
        return response
    except Exception as e:
        logging.error(f"Error serving download for file {filename}: {e}", exc_info=True)
        return f"Error: Unable to download file {filename}", 500

@app.route('/image-to-text', methods=['GET', 'POST'])
def image_to_text():
    extracted_text = None
    if request.method == 'POST':
        if 'image' not in request.files:
            flash('No image file selected.')
            return redirect(request.url)
        
        file = request.files['image']
        if file.filename == '':
            flash('No image file selected.')
            return redirect(request.url)
        
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            
            try:
                img = Image.open(filepath)
                extracted_text = pytesseract.image_to_string(img)
            except Exception as e:
                flash('Failed to extract text from image.')
                extracted_text = None
            
            # Optionally, remove the uploaded file after processing
            # os.remove(filepath)
        else:
            flash('Please upload a valid image file.')
            return redirect(request.url)
    
    return render_template('partials/image_to_text.html', extracted_text=extracted_text)

@app.route('/security')
def security():
    return render_template('security.html')

@app.route('/features')
def features():
    return render_template('features.html')

@app.route('/about')
def about():
    return render_template('about.html')

@app.route('/faq')
def faq():
    return render_template('faq.html')

@app.route('/tools')
def tools():
    return render_template('tools.html')

@app.route('/legal')
def legal():
    return render_template('legal.html')

@app.route('/contact')
def contact():
    return render_template('contact.html')

@app.route('/submit_feedback', methods=['POST'])
def submit_feedback():
    name = request.form.get('name')
    email = request.form.get('email')
    message = request.form.get('message')

    if not name or not email or not message:
        flash('All fields are required.', 'error')
        return redirect(url_for('contact'))

    # Save feedback to a file
    feedback_entry = f"Name: {name}\nEmail: {email}\nMessage: {message}\n---\n"
    feedback_file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'feedback.txt')
    try:
        with open(feedback_file_path, 'a', encoding='utf-8') as f:
            f.write(feedback_entry)
    except Exception as e:
        flash('Failed to save feedback. Please try again later.', 'error')
        return redirect(url_for('contact'))

    flash('Thank you for your feedback!', 'success')
    return redirect(url_for('contact'))

from pdf2image import convert_from_path
import tempfile
import os

from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase.pdfmetrics import stringWidth
import io

@app.route('/uploads/<path:filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/ocr_pdf', methods=['GET', 'POST'])
def ocr_pdf():
    if request.method == 'POST':
        pdf_file = request.files.get('pdf_file')
        if not pdf_file or pdf_file.filename == '':
            flash('No file selected', 'error')
            return redirect(request.url)
        if not allowed_file(pdf_file.filename) or pdf_file.filename.rsplit('.', 1)[1].lower() != 'pdf':
            flash('Please upload a valid PDF file.', 'error')
            return redirect(request.url)
        
        filename = secure_filename(pdf_file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        pdf_file.save(pdf_path)
        
        try:
            # Convert PDF pages to images
            images = convert_from_path(pdf_path)
            
            # Prepare a new PDF writer
            writer = PdfWriter()
            
            # Register a standard font
            # Use built-in Helvetica font without registering external TTF file
            # pdfmetrics.registerFont(TTFont('Helvetica', 'Helvetica.ttf'))
            
            for page_num, img in enumerate(images):
                # Perform OCR to get text and bounding boxes
                ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
                
                # Create a PDF page with the image as background
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=img.size)
                can.drawInlineImage(img, 0, 0, width=img.width, height=img.height)
                
                # Overlay invisible text
                n_boxes = len(ocr_data['level'])
                for i in range(n_boxes):
                    (x, y, w, h) = (ocr_data['left'][i], ocr_data['top'][i], ocr_data['width'][i], ocr_data['height'][i])
                    text = ocr_data['text'][i]
                    if text.strip():
                        can.setFillColorRGB(1, 1, 1, alpha=0)  # Invisible text
                        can.setFont("Helvetica", h)
                        can.drawString(x, img.height - y - h, text)
                can.save()
                
                packet.seek(0)
                new_pdf = PdfReader(packet)
                
                # Merge the image page and text overlay
                page = new_pdf.pages[0]
                writer.add_page(page)
            
            # Save the searchable PDF
            output_filename = filename.rsplit('.', 1)[0] + '_searchable.pdf'
            output_path = os.path.join(app.config['PROCESSED_FOLDER'], output_filename)
            with open(output_path, 'wb') as f_out:
                writer.write(f_out)
            
            flash('OCR processing completed successfully.', 'success')
            # Trigger auto download of the searchable PDF
            return render_template('ocr_pdf.html', download_filename=output_filename)
        except Exception as e:
            flash(f'OCR processing failed: {e}', 'error')
            return redirect(request.url)
    return render_template('ocr_pdf.html')

if __name__ == '__main__':
    app.run(debug=True)
